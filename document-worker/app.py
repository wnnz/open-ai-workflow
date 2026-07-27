import io
from pathlib import Path

from docx import Document
from fastapi import FastAPI, File, HTTPException, UploadFile, status
from openpyxl import load_workbook
from pptx import Presentation
from pydantic import BaseModel
from pypdf import PdfReader, PdfWriter
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen.canvas import Canvas

app = FastAPI(title="Open Workflow Document Worker", docs_url=None)


class PdfCreateRequest(BaseModel):
    title: str = "Generated document"
    paragraphs: list[str] = []


@app.get("/health")
def health() -> dict[str, str]:
    return {"status": "ok"}


@app.post("/inspect")
async def inspect_document(file: UploadFile = File(...)) -> dict:
    content = await file.read()
    suffix = Path(file.filename or "").suffix.lower()
    try:
        if suffix == ".xlsx":
            workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=False)
            sheets = []
            for sheet in workbook.worksheets:
                rows = [[cell.value for cell in row] for row in sheet.iter_rows(max_row=200, max_col=50)]
                sheets.append({"name": sheet.title, "rows": rows, "max_row": sheet.max_row, "max_column": sheet.max_column})
            return {"type": "excel", "sheets": sheets}
        if suffix == ".docx":
            document = Document(io.BytesIO(content))
            return {"type": "word", "paragraphs": [item.text for item in document.paragraphs], "tables": [[[cell.text for cell in row.cells] for row in table.rows] for table in document.tables]}
        if suffix == ".pptx":
            presentation = Presentation(io.BytesIO(content))
            slides = []
            for index, slide in enumerate(presentation.slides, 1):
                texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                slides.append({"number": index, "texts": texts})
            return {"type": "powerpoint", "slides": slides}
        if suffix == ".pdf":
            reader = PdfReader(io.BytesIO(content))
            return {"type": "pdf", "pages": [{"number": index + 1, "text": page.extract_text() or ""} for index, page in enumerate(reader.pages)]}
    except Exception as exc:
        raise HTTPException(status.HTTP_422_UNPROCESSABLE_ENTITY, f"Unable to parse document: {exc}") from exc
    raise HTTPException(status.HTTP_415_UNSUPPORTED_MEDIA_TYPE, "Unsupported document type")


@app.post("/pdf/create")
def create_pdf(payload: PdfCreateRequest):
    target = io.BytesIO()
    canvas = Canvas(target, pagesize=A4)
    width, height = A4
    y = height - 64
    canvas.setFont("Helvetica-Bold", 16)
    canvas.drawString(54, y, payload.title[:90])
    canvas.setFont("Helvetica", 10)
    for paragraph in payload.paragraphs:
        y -= 24
        for line in [paragraph[index:index + 100] for index in range(0, len(paragraph), 100)]:
            if y < 50:
                canvas.showPage(); y = height - 54; canvas.setFont("Helvetica", 10)
            canvas.drawString(54, y, line); y -= 14
    canvas.save()
    return {"content_base64": __import__("base64").b64encode(target.getvalue()).decode(), "mime_type": "application/pdf"}


@app.post("/pdf/merge")
async def merge_pdfs(files: list[UploadFile] = File(...)) -> dict:
    writer = PdfWriter()
    for file in files:
        reader = PdfReader(io.BytesIO(await file.read()))
        for page in reader.pages: writer.add_page(page)
    target = io.BytesIO(); writer.write(target)
    return {"content_base64": __import__("base64").b64encode(target.getvalue()).decode(), "mime_type": "application/pdf", "pages": len(writer.pages)}
